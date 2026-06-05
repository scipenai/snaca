#!/usr/bin/env python3
"""Office-format text extractor invoked by the `office-extract` skill.

Auto-detects the format from the file extension and dispatches to one of
the three OOXML Python libraries:

    .docx / .docm  → python-docx
    .xlsx / .xlsm  → openpyxl
    .pptx / .pptm  → python-pptx

Output goes to stdout as UTF-8 plain text. Errors go to stderr and the
script exits with a non-zero status so the calling Bash tool surfaces
the failure to the model. The Bash tool's 1 MB output cap is enforced
inside this script (see MAX_BYTES) so the model sees a clear truncation
marker rather than getting silently chopped at the harness boundary.
"""
from __future__ import annotations

import os
import sys

# Stay safely under Bash's 1 MiB stdout cap. The marker below adds a few
# dozen bytes; pick a round-number ceiling that leaves headroom.
MAX_BYTES = 900_000


def _fail(msg: str, code: int = 2) -> "NoReturn":  # type: ignore[name-defined]
    print(msg, file=sys.stderr)
    sys.exit(code)


def _emit(text: str) -> None:
    data = text.encode("utf-8", errors="replace")
    if len(data) > MAX_BYTES:
        truncated = data[:MAX_BYTES]
        # Snap to the nearest UTF-8 boundary so we don't print half a
        # multi-byte sequence and confuse the model.
        while truncated and (truncated[-1] & 0xC0) == 0x80:
            truncated = truncated[:-1]
        sys.stdout.buffer.write(truncated)
        sys.stdout.write(
            f"\n\n[truncated at {len(truncated)} bytes; "
            f"original was {len(data)} bytes — re-run on a smaller slice "
            f"if you need the rest]\n"
        )
    else:
        sys.stdout.buffer.write(data)
    sys.stdout.flush()


def extract_docx(path: str) -> str:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except ImportError as exc:
        _fail(
            "python-docx is not installed. "
            "Run `pip install python-docx` on the host.\n"
            f"Underlying error: {exc}"
        )
    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except ImportError as exc:
        _fail(
            "openpyxl is not installed. "
            "Run `pip install openpyxl` on the host.\n"
            f"Underlying error: {exc}"
        )
    # data_only=True so we get computed values rather than formula strings.
    wb = load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            parts.append("\t".join(cells))
        parts.append("")
    return "\n".join(parts)


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        _fail(
            "python-pptx is not installed. "
            "Run `pip install python-pptx` on the host.\n"
            f"Underlying error: {exc}"
        )
    prs = Presentation(path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"=== Slide {idx} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        parts.append(text)
        notes = getattr(slide, "notes_slide", None)
        if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
            parts.append("[notes]")
            parts.append(notes.notes_text_frame.text)
        parts.append("")
    return "\n".join(parts)


DISPATCH = {
    ".docx": extract_docx,
    ".docm": extract_docx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".pptx": extract_pptx,
    ".pptm": extract_pptx,
}


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        _fail(f"usage: {argv[0] if argv else 'extract.py'} <path-to-office-file>", code=64)
    path = argv[1]
    if not os.path.isfile(path):
        _fail(f"not a regular file: {path}", code=66)
    ext = os.path.splitext(path)[1].lower()
    handler = DISPATCH.get(ext)
    if handler is None:
        _fail(
            f"unsupported extension {ext!r}; expected one of: "
            + ", ".join(sorted(DISPATCH))
        )
    try:
        text = handler(path)
    except Exception as exc:  # extractor library blew up
        _fail(f"extraction failed for {path}: {exc}", code=70)
    _emit(text)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
